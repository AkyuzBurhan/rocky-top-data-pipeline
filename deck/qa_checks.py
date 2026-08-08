"""
Programmatic QA for the built deck against the binding design system and
plan.md. Renders are inspected separately; this checks what XML can prove.

Usage: python deck/qa_checks.py
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
PPTX = HERE / "out" / "BZAN545_Final_Deck.pptx"
# canonical home is docs/plan.md (close-out); root plan.md is the pre-move
# fallback so the check runs in either state
_PLAN_CANDIDATES = (HERE.parent / "docs" / "plan.md", HERE.parent / "plan.md")
PLAN = next((p for p in _PLAN_CANDIDATES if p.exists()), _PLAN_CANDIDATES[0])

FONT_WHITELIST = {"IBM Plex Sans", "IBM Plex Sans SmBld", "IBM Plex Sans Medm",
                  "IBM Plex Mono", "IBM Plex Mono SmBld"}
ORANGE = "FF8200"
# expected raw orange XML marks per slide (shapes + runs); one design
# element can serialize as several runs. Slides 2 and 4 carry the legend's
# open-dot key next to the incident dot (Connor item 5); slide 11 marks the
# "14 of 18" phrase = 3 runs; slide 14 (A2) marks the 1mm row = 3 runs.
EXPECTED_ORANGE = [1, 2, 1, 2, 1, 1, 1, 1, 1, 1, 3, 1, 1, 3, 1]
TITLE_X_EMU = 548640      # 0.6 in
TITLE_Y_EMU = 502920      # 0.55 in
ALLOWED_LINE_W = {9525, 12700}          # 0.75pt hairlines, 1pt strokes
NUM_RE = re.compile(r"\$?\d[\d,.\-]*\d%?|\$?\d%?")
# On-slide numeric tokens not present in plan.md, each cited in the readout.
NUM_ALLOW = {
    "157.12",     # data/raw/orders_2026-08-05.csv:9 ($157.12 string)
    "07-23",      # ingestion_log.csv: source_date_found=2026-07-23
    "07-07",      # window start, plan.md "07-07 -> 08-07"
    "2026",       # cover date line "August 2026"
    "545",        # course code, plan.md title
    "101-103",    # load_raw.py citation in notes-adjacent A3/caption text
}
# poison-fixture tokens (A3): date, code commit, file/flag figures
_meta = Path(__file__).resolve().parent / "out/fixture/poison_fixture_meta.json"
if _meta.exists():
    import json
    _m = json.loads(_meta.read_text(encoding="utf-8"))
    NUM_ALLOW |= {_m["commit"], _m["date"], "2026-08-08", "08-08",
                  "503-504", "565-566"}
    # a sha like 1c4884f tokenizes as its digit runs; allow those too
    NUM_ALLOW |= set(re.findall(r"\d[\d,.\-]*\d|\d", _m["commit"]))
FOOTER_LABELS = {"1", "2", "3", "4", "5", "6", "7a", "7b", "8", "9", "10",
                 "A1", "A2", "A3"}


def iter_text_shapes(shapes):
    for shp in shapes:
        if shp.shape_type == 6:                     # group
            yield from iter_text_shapes(shp.shapes)
        yield shp


def shape_texts(slide):
    """(text, is_footer) per text-bearing shape, tables included."""
    out = []
    for shp in iter_text_shapes(slide.shapes):
        if shp.has_text_frame:
            txt = "\n".join(p.text for p in shp.text_frame.paragraphs)
            is_footer = (shp.top and Emu(shp.top).inches > 7.0
                         and txt.strip() in FOOTER_LABELS)
            out.append((txt, is_footer))
        if getattr(shp, "has_table", False) and shp.has_table:
            for row in shp.table.rows:
                for cell in row.cells:
                    out.append((cell.text, False))
    return out


def run_fonts(slide):
    missing, bad = [], []
    for shp in iter_text_shapes(slide.shapes):
        frames = []
        if shp.has_text_frame:
            frames.append(shp.text_frame)
        if getattr(shp, "has_table", False) and shp.has_table:
            frames += [c.text_frame for r in shp.table.rows for c in r.cells]
        for tf in frames:
            for p in tf.paragraphs:
                for r in p.runs:
                    name = r.font.name
                    if not r.text.strip():
                        continue
                    if name is None:
                        missing.append(r.text[:30])
                    elif name not in FONT_WHITELIST:
                        bad.append((name, r.text[:30]))
    return missing, bad


def orange_count(slide):
    n = 0
    xml = slide._element.xml
    # shape lines/fills and text runs both serialize as srgbClr FF8200
    n += xml.count(ORANGE)
    return n


def check(deck_path=PPTX):
    prs = Presentation(deck_path)
    plan_text = PLAN.read_text(encoding="utf-8")
    plan_tokens = set(NUM_RE.findall(plan_text))
    failures, notes = [], []

    for idx, slide in enumerate(prs.slides, start=1):
        xml = slide._element.xml
        label = f"slide {idx}"

        # --- ban sweep ---
        all_text = "\n".join(t for t, _ in shape_texts(slide))
        notes_text = (slide.notes_slide.notes_text_frame.text
                      if slide.has_notes_slide else "")
        for banned, where in [("—", "em dash"), ("•", "bullet char"),
                              ("[surname]", "surname placeholder")]:
            if banned in all_text:
                failures.append(f"{label}: {where} in slide text")
            if banned in notes_text:
                failures.append(f"{label}: {where} in notes")
        # A3 (slide 15) carries the executed poison-fixture log capture as
        # a raster by Connor's item 24; images stay banned everywhere else
        for frag, why in [("4472C4", "Office blue"), ("gradFill", "gradient"),
                          ('prstGeom prst="roundRect"', "rounded corners"),
                          ("outerShdw", "drop shadow")]:
            if frag in xml:
                failures.append(f"{label}: banned {why} in XML")
        if "<pic:pic" in xml and idx != 15:
            failures.append(f"{label}: banned image in XML")

        # --- fonts ---
        missing, bad = run_fonts(slide)
        if missing:
            failures.append(f"{label}: {len(missing)} runs with no font: "
                            f"{missing[:3]}")
        if bad:
            failures.append(f"{label}: non-Plex fonts {bad[:3]}")

        # --- orange ---
        n_orange = orange_count(slide)
        want = EXPECTED_ORANGE[idx - 1]
        if n_orange != want:
            failures.append(f"{label}: orange marks {n_orange} != "
                            f"expected {want}")

        # --- title block position (content slides only) ---
        if idx > 1:
            tops = [(shp.left, shp.top) for shp in slide.shapes
                    if shp.has_text_frame and shp.top is not None
                    and abs(shp.top - TITLE_Y_EMU) < 2000]
            if not any(abs(left - TITLE_X_EMU) < 2000 for left, _ in tops):
                failures.append(f"{label}: no title block at 0.6/0.55 in")

        # --- stroke widths ---
        for m in re.finditer(r'<a:ln[ >][^>]*w="(\d+)"', xml):
            w = int(m.group(1))
            if w not in ALLOWED_LINE_W:
                failures.append(f"{label}: stroke width {w / 12700:.2f}pt")

        # --- incident glyph geometry: every ellipse is a circle ---
        for shp in iter_text_shapes(slide.shapes):
            if 'prst="ellipse"' in shp._element.xml:
                if shp.width != shp.height:
                    failures.append(f"{label}: oval {Emu(shp.width).inches}"
                                    f"x{Emu(shp.height).inches} not circular")

        # --- number diff vs plan.md ---
        strange = set()
        for text, is_footer in shape_texts(slide):
            if is_footer:
                continue
            for tok in NUM_RE.findall(text):
                # units live in the header, not cells: try the %-suffixed
                # and $-stripped variants a plan.md figure may carry
                variants = {tok, tok + "%", tok.lstrip("$")}
                if variants & plan_tokens or variants & NUM_ALLOW:
                    continue
                strange.add(tok)
        if strange:
            failures.append(f"{label}: numbers not in plan.md/allowlist: "
                            f"{sorted(strange)}")

    # incident strips: 6 dots slide 2 (+2 legend), 6 dots slide 5 mini
    for idx, want_ovals in [(2, 8), (4, 8), (5, 6)]:
        n = sum(1 for shp in iter_text_shapes(prs.slides[idx - 1].shapes)
                if 'prst="ellipse"' in shp._element.xml)
        if n != want_ovals:
            failures.append(f"slide {idx}: {n} ovals, expected {want_ovals}")

    print(f"QA against {deck_path.name}: {len(failures)} failure(s)")
    for f in failures:
        print("  FAIL", f)
    for n in notes:
        print("  note", n)
    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(check())
