"""
Phase-1 specimen slides: title slide, slide 3 (incident table), slide 5
(matcher funnel + mechanism box). Content transcribed from ../plan.md.

Slide-3 cell text for failure / detection layer / outcome is composed from
docs/DECISIONS.md incident entries (plan.md names only the six rows); every
figure is plan.md verbatim. Flagged for review at the specimen gate.

Usage:  python deck/build_specimens.py
Output: deck/out/specimens/specimens.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

import design as d

HERE = Path(__file__).resolve().parent
OUT = HERE / "out" / "specimens"


def new_slide(prs, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    d.set_background(slide)
    if footer is not None:
        d.add_footer(slide, footer)
    return slide


# --- title slide ----------------------------------------------------------

def title_slide(prs):
    s = new_slide(prs)
    d.add_text(s, d.MARGIN, 2.45, d.CONTENT_W, 0.4,
               [[("BZAN 545 · Rocky Top Outfitters data pipeline",
                  d.SANS, 14, d.GRAY)]])
    d.add_text(s, d.MARGIN, 2.95, d.CONTENT_W, 1.8,
               [[("A process that notices", d.SANS_SB, 40, d.INK)],
                [("when it's wrong", d.SANS_SB, 40, d.INK),
                 (".", d.SANS_SB, 40, d.ORANGE)]],
               space_after=0, line_spacing=1.08)
    d.add_text(s, d.MARGIN, 7.08, 4.0, 0.3,
               [[("August 2026", d.SANS, d.FOOT_SIZE, d.GRAY)]])
    d.add_notes(s, [
        "Connor · pre-start",
        "- Open the dashboard before class: Streamlit Community Cloud "
        "sleeps after 12 hours; slide 7's chart must not cold-start",
        "- Win+P -> Extend (not Duplicate) for presenter view",
        "- Deck locked the night before; randomized order, slide 1 works "
        "cold at 10:00 AM",
    ])
    return s


# --- slide 3: incident table ----------------------------------------------

INCIDENTS = [
    # date, failure, detection layer, outcome, missed
    ("07-24", "Stale file re-served: 07-23 data again",
     "Capture — stale-file check",
     "Flagged; 140 duplicate rows quarantined", False),
    ("07-28", "Product-ID migration, P#### → NP####",
     "Quality flag — new-ID schema",
     "Crosswalk built; legacy keys retained", False),
    ("08-03", "Empty file: header row, zero data rows",
     "Capture — empty-file check",
     "Logged empty; nothing invented", False),
    ("08-05", 'Prices arrived as "$157.12" strings',
     "Audit grep; no check fired",
     "$56,970.09 recovered from raw; parser fixed", True),
    ("08-06", "Source URL returned 404",
     "Capture — HTTP status",
     "Logged failed; clean recovery 08-07", False),
    ("08-07", "Column reorder in orders CSV",
     "Design — header-based reads",
     "Zero impact; loaded correctly", False),
]

FOOTNOTE_3 = ("144 rejected = 140 stale-file dupes + 4 within-file dupes "
              "07-16; checks earn their keep on ordinary days too")


def slide_3(prs):
    s = new_slide(prs, footer="3")
    d.add_title(s, "Six incidents in 32 days; five caught at the layer "
                   "where they occurred")

    tbl_x, tbl_y = 0.95, 2.05
    header_h, row_h = 0.38, 0.55
    col_w = [1.15, 4.15, 3.0, 3.48]
    cells = [["Date", "Failure", "Detection layer", "Outcome"]]
    for date, failure, layer, outcome, _ in INCIDENTS:
        cells.append([[(date, d.MONO, d.BODY_SIZE, d.INK)],
                      failure, layer, outcome])
    d.tufte_table(s, tbl_x, tbl_y, col_w,
                  [header_h] + [row_h] * len(INCIDENTS), cells)

    # incident glyphs in a gutter left of the table, one per row
    for i, (_, _, _, _, missed) in enumerate(INCIDENTS):
        cy = tbl_y + header_h + row_h * i + row_h / 2
        d.incident_dot(s, 0.72, cy, 0.14, missed=missed)

    d.add_text(s, tbl_x, 6.55, d.CONTENT_W - (tbl_x - d.MARGIN), 0.4,
               [[(FOOTNOTE_3, d.SANS, d.FOOT_SIZE, d.GRAY)]])

    d.add_notes(s, [
        "Burhan, 2:00",
        "- Six incidents, 32 days; five caught at the layer where they "
        "occurred",
        "- Walk the layers: capture caught stale (07-24), empty (08-03), "
        "404 (08-06); quality flag caught the ID migration (07-28); "
        "header-based reads absorbed the column reorder (08-07)",
        "- 08-05 is the open dot: the one that got through the pipeline",
        "- Footnote beat: 144 rejected = 140 stale-file dupes + 4 "
        "within-file dupes (07-16); checks earn their keep on ordinary "
        "days too",
        "- Handoff: 'the one that got through is Jack's story'",
        "Cut order 3 (only under real pressure): compress this table onto "
        "slide 2's diagram, -1.0; costs Burhan airtime",
    ])
    return s


# --- slide 5: matcher funnel + mechanism box ------------------------------

MECHANISM_1 = ("Two thresholds, two jobs. Composite (0.6·name + "
               "0.2·subclass + 0.2·price) ≥ 0.60 decides accept vs reject. "
               "name_sim ≥ 0.85 decides high vs medium confidence.")
MECHANISM_2 = "Review queue = everything not high (14 medium + 4 low = 18)."


def _node(s, x, y, w, h, number, label, num_size=20):
    d.stroke_box(s, x, y, w, h)
    d.add_text(s, x + 0.14, y, w - 0.28, h,
               [[(number, d.MONO_SB, num_size, d.INK)],
                [(label, d.SANS, d.LABEL_SIZE, d.GRAY)]],
               anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)


def _chip(s, x, y, w, h, number, label):
    d.stroke_box(s, x, y, w, h)
    d.add_text(s, x + 0.12, y, w - 0.24, h,
               [[(number, d.MONO_SB, 14, d.INK),
                 ("  " + label, d.SANS, d.LABEL_SIZE, d.INK)]],
               anchor=MSO_ANCHOR.MIDDLE, wrap=False, space_after=0)


def slide_5(prs):
    s = new_slide(prs, footer="5")
    d.add_title(s, "The matcher grades its own confidence: 76 of 80 mapped, "
                   "and all 18 it's less than sure of are flagged with a "
                   "reason")

    # source
    _node(s, 0.6, 3.0, 1.7, 1.2, "80", "legacy products", num_size=28)

    # tiers
    tiers = [
        (2.0, "51", "exact name"),
        (3.2, "25", "fuzzy within block"),
        (4.4, "4", "no candidate"),
    ]
    for y, num, label in tiers:
        _node(s, 3.0, y, 2.9, 0.9, num, label)
        d.connector(s, 2.3, 3.6, 3.0, y + 0.45, elbow=True)

    # confidence-tier chips
    d.add_text(s, 6.6, 1.68, 2.7, 0.3,
               [[("confidence tier", d.SANS, d.LABEL_SIZE, d.GRAY)]])
    chips = [
        (2.04, "51", "→  high", 2.45),
        (2.74, "11", "→  high", 3.65),
        (3.39, "10", "→  matched · medium", 3.65),
        (4.04, "4", "→  possible · medium", 3.65),
        (4.74, "4", "→  low", 4.85),
    ]
    for y, num, label, src_cy in chips:
        _chip(s, 6.6, y, 2.7, 0.52, num, label)
        d.connector(s, 5.9, src_cy, 6.6, y + 0.26, elbow=True)

    # key figures, display treatment
    d.add_text(s, 10.0, 2.1, 2.73, 0.85,
               [[("76", d.MONO, 44, d.INK), (" / 80", d.MONO, 20, d.GRAY)]],
               wrap=False, space_after=0, line_spacing=1.0)
    d.add_text(s, 10.0, 2.95, 2.73, 0.35,
               [[("legacy products mapped", d.SANS, d.LABEL_SIZE, d.GRAY)]])
    d.add_text(s, 10.0, 3.9, 2.73, 0.85,
               [[("18", d.MONO, 44, d.ORANGE)]],
               wrap=False, space_after=0, line_spacing=1.0)
    d.add_text(s, 10.0, 4.75, 2.73, 0.35,
               [[("flagged for review, with a reason",
                  d.SANS, d.LABEL_SIZE, d.GRAY)]])

    # mechanism box
    d.stroke_box(s, 0.6, 5.55, d.CONTENT_W, 1.25, fill=d.TINT5, stroke=None)
    d.add_text(s, 0.85, 5.55, d.CONTENT_W - 0.5, 1.25,
               [("mix", MECHANISM_1, d.BODY_SIZE, d.INK),
                ("mix", MECHANISM_2, d.BODY_SIZE, d.INK)],
               anchor=MSO_ANCHOR.MIDDLE, space_after=4)

    d.add_notes(s, [
        "Burhan, 2:30",
        "- Funnel: 80 legacy -> 51 exact high · 25 fuzzy (11 high, 10 "
        "matched-medium, 4 possible-medium) · 4 low",
        "- Two thresholds, two jobs: composite >= 0.60 accept vs reject; "
        "name_sim >= 0.85 high vs medium",
        "- Review queue = everything not high: 14 medium + 4 low = 18",
        "- P1076 exhibit: weakest accept cleared the floor by 0.068 and "
        "got flagged; thin acceptance plus automatic review is the whole "
        "design in one product",
        "- Three-threshold motif: MIN_SCORE, STRONG_NAME, and the rain cut "
        "that killed the revenue claim; one sentence on 7b, don't belabor",
        "- Handoff: 'the flag is advisory; back to that in limitations'",
        "Q2 (not Burhan) — P1076 matched at 0.50 name similarity against a "
        "0.60 floor: 'The floor applies to the composite, not the name. "
        "P1076's name scored 0.50, but subclass matched exactly and price "
        "closeness was 0.841, putting the composite at 0.668. Its block "
        "had one candidate. And because name similarity was under the 0.85 "
        "confidence threshold, it came out medium and landed in the review "
        "queue: the system flagged its own weakest accept.'",
        "Q3 (anyone) — 72 matched but 18 need review, which is it? "
        "Orthogonal axes: status = decision made, confidence = "
        "name-evidence strength. 10 of the 72 are medium: committed "
        "downstream, queued for audit.",
    ])
    return s


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = d.SLIDE_W_EMU
    prs.slide_height = d.SLIDE_H_EMU
    title_slide(prs)
    slide_3(prs)
    slide_5(prs)
    path = OUT / "specimens.pptx"
    prs.save(path)
    print(f"wrote {path}")


if __name__ == "__main__":
    main()
