"""
Build the full BZAN 545 deck: cover + slides 1-6, 7a, 7b, 8-10 + appendix
A1-A3. Content comes from content.py (transcribed from ../plan.md); layout
and chrome from design.py.

Usage:  python deck/build_deck.py
Output: deck/out/BZAN545_Final_Deck.pptx
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

import content as c
import design as d

HERE = Path(__file__).resolve().parent
OUT = HERE / "out"

# Slide 1's 158-char title: 22pt floor pre-authorized on that slide only.
S1_TITLE_SIZE = 22


def new_slide(prs, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    d.set_background(slide)
    if footer is not None:
        d.add_footer(slide, footer)
    return slide


def legend(slide, x, y, size=0.10, gap=0.16):
    """Two-item incident-glyph legend, identical wording on slides 1 and 3.
    The open key is stroked orange like the glyph it explains; QA counts it
    as a legend exemption."""
    d.incident_dot(slide, x + size / 2, y, size, missed=False)
    d.add_text(slide, x + size + 0.08, y - 0.09, 2.2, 0.25,
               [[(c.LEGEND_CAUGHT, d.SANS, d.FOOT_SIZE, d.GRAY)]],
               wrap=False, space_after=0)
    x2 = x + size + 0.08 + 1.35 + gap
    d.incident_dot(slide, x2 + size / 2, y, size, missed=True)
    d.add_text(slide, x2 + size + 0.08, y - 0.09, 2.2, 0.25,
               [[(c.LEGEND_MISSED, d.SANS, d.FOOT_SIZE, d.GRAY)]],
               wrap=False, space_after=0)


# --- cover ----------------------------------------------------------------

def cover(prs):
    s = new_slide(prs)
    d.add_text(s, d.MARGIN, 2.45, d.CONTENT_W, 0.4,
               [[(c.COVER_KICKER, d.SANS, 14, d.GRAY)]])
    d.add_text(s, d.MARGIN, 2.95, d.CONTENT_W, 1.8,
               [[(c.COVER_TITLE_1, d.SANS_SB, 40, d.INK)],
                [(c.COVER_TITLE_2, d.SANS_SB, 40, d.INK),
                 (".", d.SANS_SB, 40, d.ORANGE)]],
               space_after=0, line_spacing=1.08)
    d.add_text(s, d.MARGIN, 6.72, d.CONTENT_W, 0.3,
               [[(c.COVER_ROSTER, d.SANS, 12, d.GRAY)]])
    d.add_text(s, d.MARGIN, 7.08, 4.0, 0.3,
               [[(c.COVER_DATE, d.SANS, d.FOOT_SIZE, d.GRAY)]])
    d.add_notes(s, c.COVER_NOTES)


# --- slide 1 --------------------------------------------------------------

def slide_1(prs):
    s = new_slide(prs, footer="1")
    box = d.add_title(s, c.S1_TITLE)
    for p in box.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(S1_TITLE_SIZE)

    axis_y = 4.15
    x0, w = d.MARGIN, d.CONTENT_W
    d.add_text(s, x0, 3.30, 5.0, 0.3,
               [[(c.WINDOW_LABEL, d.SANS, d.LABEL_SIZE, d.GRAY)]])
    d.hairline(s, x0, axis_y, w)
    # dots proportional to day offset in the 0..31 window
    label_rows = {"08-05": 1, "08-07": 1}      # stagger the tight cluster
    for date, day in c.S1_INCIDENTS:
        cx = x0 + (day / 31) * w
        d.incident_dot(s, cx, axis_y, 0.22, missed=(date == c.S1_MISSED))
        row = label_rows.get(date, 0)
        ly = axis_y + 0.24 + row * 0.26
        if date == "08-07":                     # keep inside right margin
            d.add_text(s, cx - 0.9, ly, 0.9, 0.25,
                       [[(date, d.MONO, d.FOOT_SIZE, d.GRAY)]],
                       align=PP_ALIGN.RIGHT, wrap=False, space_after=0)
        else:
            d.add_text(s, cx - 0.45, ly, 0.9, 0.25,
                       [[(date, d.MONO, d.FOOT_SIZE, d.GRAY)]],
                       align=PP_ALIGN.CENTER, wrap=False, space_after=0)
    legend(s, x0, 5.45)
    d.add_notes(s, c.S1_NOTES + [c.WINDOW_NOTE])


# --- slide 2 --------------------------------------------------------------

def slide_2(prs):
    s = new_slide(prs, footer="2")
    d.add_title(s, c.S2_TITLE)

    bw, bh, y = 1.9, 0.85, 2.7
    gap = (d.CONTENT_W - 5 * bw) / 4
    for i, stage in enumerate(c.S2_STAGES):
        x = d.MARGIN + i * (bw + gap)
        stroke = d.ORANGE if stage == "raw" else d.INK
        d.stroke_box(s, x, y, bw, bh, stroke=stroke)
        d.add_text(s, x, y, bw, bh, [[(stage, d.SANS, 14, d.INK)]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                   space_after=0)
        if i > 0:
            d.connector(s, x - gap + 0.02, y + bh / 2, x - 0.02, y + bh / 2)
        tag = {"raw": c.S2_TAG_RAW, "transform": c.S2_TAG_TRANSFORM}.get(stage)
        if tag:
            d.add_text(s, x - 0.3, y + bh + 0.12, bw + 0.6, 0.3,
                       [[(tag, d.SANS, d.LABEL_SIZE, d.GRAY)]],
                       align=PP_ALIGN.CENTER, wrap=False, space_after=0)

    d.stroke_box(s, 7.4, 5.5, 5.33, 1.3, fill=d.TINT5, stroke=None)
    d.add_text(s, 7.65, 5.5, 4.83, 1.3,
               [("mix", c.S2_CORNER, d.LABEL_SIZE, d.INK)],
               anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    d.add_notes(s, c.S2_NOTES)


# --- slide 3 --------------------------------------------------------------

def slide_3(prs):
    s = new_slide(prs, footer="3")
    d.add_title(s, c.S3_TITLE)

    tbl_x, tbl_y = 0.95, 2.05
    header_h, row_h = 0.38, 0.55
    col_w = [1.15, 4.15, 3.0, 3.48]
    cells = [["Date", "Failure", "Detection layer", "Outcome"]]
    for date, failure, layer, outcome, _ in c.S3_ROWS:
        cells.append([[(date, d.MONO, d.BODY_SIZE, d.INK)],
                      failure, layer, outcome])
    d.tufte_table(s, tbl_x, tbl_y, col_w,
                  [header_h] + [row_h] * len(c.S3_ROWS), cells)
    for i, (_, _, _, _, missed) in enumerate(c.S3_ROWS):
        cy = tbl_y + header_h + row_h * i + row_h / 2
        d.incident_dot(s, 0.72, cy, 0.14, missed=missed)

    legend(s, d.MARGIN, 6.08)
    d.add_text(s, tbl_x, 6.48, d.CONTENT_W - (tbl_x - d.MARGIN), 0.3,
               [[(c.S3_FOOTNOTE, d.SANS, d.FOOT_SIZE, d.GRAY)]])
    d.add_text(s, tbl_x, 6.80, d.CONTENT_W - (tbl_x - d.MARGIN), 0.3,
               [[(c.WINDOW_FOOTLINE, d.SANS, d.FOOT_SIZE, d.GRAY)]])
    notes = list(c.S3_NOTES)
    notes.insert(len(notes) - 1, c.WINDOW_NOTE)
    d.add_notes(s, notes)


# --- slide 4 --------------------------------------------------------------

def slide_4(prs):
    s = new_slide(prs, footer="4")
    box = d.add_title(s, c.S4_TITLE)
    box.width = d.Inches(10.4)          # reserve top-right for the mini strip

    # mini incident strip: sequence recurrence, open dot orange (slide's one
    # orange element), dot geometry identical to slide 1 at 0.15in
    d.incident_strip(s, 11.28, 0.78, 1.45, 0.15,
                     missed_index=[i for i, (dt, _) in
                                   enumerate(c.S1_INCIDENTS)
                                   if dt == c.S1_MISSED][0])

    # beat 1: the three passed checks
    d.add_text(s, d.MARGIN, 2.35, 2.6, 0.3,
               [[(c.S4_CHECKS_LABEL, d.SANS, d.LABEL_SIZE, d.GRAY)]])
    for i, chk in enumerate(c.S4_CHECKS):
        y = 2.7 + i * 0.62
        d.stroke_box(s, d.MARGIN, y, 2.6, 0.48)
        d.add_text(s, d.MARGIN + 0.12, y, 2.36, 0.48,
                   [[(chk, d.SANS, d.LABEL_SIZE, d.INK)]],
                   anchor=MSO_ANCHOR.MIDDLE, wrap=False, space_after=0)
    d.connector(s, 3.2, 3.55, 3.75, 3.55)

    # beat 2: the miss, said plainly
    d.stroke_box(s, 3.75, 2.7, 3.3, 1.72, fill=d.TINT8, stroke=None)
    d.add_text(s, 4.0, 2.7, 2.8, 1.72,
               [[(c.S4_MISS, d.SANS, d.BODY_SIZE, d.INK)]],
               anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    d.connector(s, 7.05, 3.55, 7.6, 3.55)

    # beat 3: the recovery figure
    d.add_text(s, 7.6, 2.95, 3.1, 0.85,
               [[(c.S4_AMOUNT, d.MONO, 40, d.INK)]],
               wrap=False, space_after=0, line_spacing=1.0)
    d.add_text(s, 7.6, 3.78, 3.1, 0.3,
               [[(c.S4_AMOUNT_LABEL, d.SANS, d.LABEL_SIZE, d.GRAY)]],
               wrap=False, space_after=0)
    d.connector(s, 10.75, 3.55, 11.15, 3.55)

    # beat 4: the fix
    d.stroke_box(s, 11.15, 3.13, 1.58, 0.85)
    d.add_text(s, 11.27, 3.13, 1.34, 0.85,
               [[(c.S4_FIX, d.SANS, d.LABEL_SIZE, d.INK)]],
               anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    d.add_notes(s, c.S4_NOTES)


# --- slide 5 --------------------------------------------------------------

def _node(s, x, y, w, h, number, label, num_size=20):
    d.stroke_box(s, x, y, w, h)
    d.add_text(s, x + 0.14, y, w - 0.28, h,
               [[(number, d.MONO_SB, num_size, d.INK)],
                [(label, d.SANS, d.LABEL_SIZE, d.GRAY)]],
               anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)


def _chip(s, x, y, w, h, number, label):
    d.stroke_box(s, x, y, w, h)
    d.add_text(s, x + 0.12, y, w - 0.24, h,
               [[(number, d.MONO_SB, 14, d.INK),
                 ("  " + label, d.SANS, d.LABEL_SIZE, d.INK)]],
               anchor=MSO_ANCHOR.MIDDLE, wrap=False, space_after=0)


def slide_5(prs):
    s = new_slide(prs, footer="5")
    d.add_title(s, c.S5_TITLE)

    _node(s, 0.6, 3.0, 1.7, 1.2, c.S5_SOURCE[0], c.S5_SOURCE[1], num_size=28)
    tier_ys = [2.0, 3.2, 4.4]
    for (num, label), y in zip(c.S5_TIERS, tier_ys):
        _node(s, 3.0, y, 2.9, 0.9, num, label)
        d.connector(s, 2.3, 3.6, 3.0, y + 0.45, elbow=True)

    d.add_text(s, 6.6, 1.68, 2.7, 0.3,
               [[(c.S5_CHIP_CAPTION, d.SANS, d.LABEL_SIZE, d.GRAY)]])
    chip_ys = [2.04, 2.74, 3.39, 4.04, 4.74]
    chip_src = [2.45, 3.65, 3.65, 3.65, 4.85]
    for (num, label), y, src_cy in zip(c.S5_CHIPS, chip_ys, chip_src):
        _chip(s, 6.6, y, 2.7, 0.52, num, label)
        d.connector(s, 5.9, src_cy, 6.6, y + 0.26, elbow=True)

    # two ink brackets at 0.75pt (item 21: hairline gray washes out at
    # projection distance), never orange: mapped (76) spans the first four
    # chips, review (18) spans the last three; they overlap on the middle
    # two by design (76 + 18 = 94 on an 80-product slide), and the 14
    # shared rows are labeled between the two stats
    xa, xb = 9.48, 9.72
    d.hairline(s, 9.36, 2.30, xa - 9.36, color=d.INK)   # bracket A stubs
    d.hairline(s, 9.36, 4.30, xa - 9.36, color=d.INK)
    d.connector(s, xa, 2.30, xa, 4.30, arrow=False, weight=0.75,
                color=d.INK)
    d.hairline(s, xa, 2.48, 9.95 - xa, color=d.INK)     # meets the 76
    d.hairline(s, 9.54, 3.65, xb - 9.54, color=d.INK)   # bracket B stubs
    d.hairline(s, 9.54, 5.00, xb - 9.54, color=d.INK)
    d.connector(s, xb, 3.65, xb, 5.00, arrow=False, weight=0.75,
                color=d.INK)
    d.hairline(s, xb, 4.28, 9.95 - xb, color=d.INK)     # meets the 18
    d.add_text(s, 10.0, 3.42, 2.73, 0.3,
               [[(c.S5_OVERLAP, d.SANS, d.FOOT_SIZE, d.GRAY)]])

    d.add_text(s, 10.0, 2.1, 2.73, 0.85,
               [[(c.S5_STAT_MAPPED[0], d.MONO, 44, d.INK),
                 (c.S5_STAT_MAPPED[1], d.MONO, 20, d.GRAY)]],
               wrap=False, space_after=0, line_spacing=1.0)
    d.add_text(s, 10.0, 2.95, 2.73, 0.35,
               [[(c.S5_STAT_MAPPED[2], d.SANS, d.LABEL_SIZE, d.GRAY)]])
    d.add_text(s, 10.0, 3.9, 2.73, 0.85,
               [[(c.S5_STAT_REVIEW[0], d.MONO, 44, d.ORANGE)]],
               wrap=False, space_after=0, line_spacing=1.0)
    d.add_text(s, 10.0, 4.75, 2.73, 0.35,
               [[(c.S5_STAT_REVIEW[1], d.SANS, d.LABEL_SIZE, d.GRAY)]])

    d.stroke_box(s, 0.6, 5.55, d.CONTENT_W, 1.25, fill=d.TINT5, stroke=None)
    d.add_text(s, 0.85, 5.55, d.CONTENT_W - 0.5, 1.25,
               [("mix", c.S5_MECHANISM_1, d.BODY_SIZE, d.INK),
                ("mix", c.S5_MECHANISM_2, d.BODY_SIZE, d.INK)],
               anchor=MSO_ANCHOR.MIDDLE, space_after=4)
    d.add_notes(s, c.S5_NOTES)


# --- slide 6 --------------------------------------------------------------

def slide_6(prs):
    s = new_slide(prs, footer="6")
    d.add_title(s, c.S6_TITLE)

    # weather block: the 232 independent observations
    d.stroke_box(s, 0.6, 2.7, 3.3, 1.9)
    d.add_text(s, 0.85, 2.7, 2.8, 1.9,
               [[(c.S6_LEFT[0], d.MONO_SB, 36, d.ORANGE)],
                [(c.S6_LEFT[1], d.SANS, d.LABEL_SIZE, d.GRAY)],
                [(c.S6_LEFT[2], d.SANS, d.LABEL_SIZE, d.GRAY)]],
               anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.05)

    # sales block: 1,367 store-day rows sharing them
    d.stroke_box(s, 9.43, 2.7, 3.3, 1.9)
    d.add_text(s, 9.68, 2.7, 2.8, 1.9,
               [[(c.S6_RIGHT[0], d.MONO_SB, 36, d.INK)],
                [(c.S6_RIGHT[1], d.SANS, d.LABEL_SIZE, d.GRAY)]],
               anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.05)

    # join fan: one observation feeds many rows
    d.add_text(s, 4.4, 2.55, 4.6, 0.3,
               [[(c.S6_JOIN_LABEL, d.SANS, d.LABEL_SIZE, d.GRAY)]],
               align=PP_ALIGN.CENTER)
    for dy in (-0.55, -0.275, 0.0, 0.275, 0.55):
        d.connector(s, 3.9, 3.65, 9.43, 3.65 + dy)

    # NULL tail from archive lag
    d.stroke_box(s, 9.43, 4.75, 3.3, 0.42, fill=d.TINT8, stroke=None)
    d.add_text(s, 9.68, 4.75, 2.8, 0.42,
               [[(c.S6_NULL_TAIL, d.SANS, d.FOOT_SIZE, d.GRAY)]],
               anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    d.add_notes(s, c.S6_NOTES)


# --- slides 7a / 7b -------------------------------------------------------

def slide_7a(prs):
    s = new_slide(prs, footer="7a")
    d.add_title(s, c.S7A_TITLE)

    col_w = [2.6, 1.7, 1.7]
    cells = [c.S7A_TABLE_HEADER]
    for ch, dry, rain in c.S7A_TABLE:
        cells.append([[(ch, d.MONO, d.BODY_SIZE, d.INK)],
                      [(dry, d.MONO, d.BODY_SIZE, d.INK)],
                      [(rain, d.MONO, d.BODY_SIZE, d.INK)]])
    d.tufte_table(s, 0.6, 2.35, col_w, [0.38] + [0.52] * 3, cells,
                  align=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])

    d.add_text(s, 8.1, 2.55, 4.6, 1.0,
               [[(c.S7A_STAT[0], d.MONO, 48, d.ORANGE)]],
               wrap=False, space_after=0, line_spacing=1.0)
    d.add_text(s, 8.1, 3.5, 4.6, 0.35,
               [[(c.S7A_STAT[1], d.SANS, d.LABEL_SIZE, d.GRAY)]])

    d.add_text(s, 0.6, 5.15, 8.0, 0.35,
               [("mix", c.S7A_CAVEAT_1, d.BODY_SIZE, d.INK)])
    d.add_text(s, 0.6, 5.55, 10.0, 0.35,
               [("mix", c.S7A_CAVEAT_2, d.LABEL_SIZE, d.GRAY)])
    d.add_notes(s, c.S7A_NOTES)


def _threshold_table(s, x, y, header, rows, orange_row=None):
    cells = [header]
    for i, (mm, lift, agree) in enumerate(rows):
        color = d.ORANGE if i == orange_row else d.INK
        cells.append([[(mm, d.MONO, d.BODY_SIZE, color)],
                      [(lift, d.MONO, d.BODY_SIZE, color)],
                      [(agree, d.MONO, d.BODY_SIZE, color)]])
    d.tufte_table(s, x, y, [1.9, 2.1, 2.0], [0.38] + [0.5] * len(rows),
                  cells,
                  align=[PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])


def slide_7b(prs):
    s = new_slide(prs, footer="7b")
    d.add_title(s, c.S7B_TITLE)

    _threshold_table(s, 0.6, 2.35, c.S7B_TABLE_HEADER, c.S7B_TABLE)

    d.add_text(s, 7.6, 2.45, 5.13, 0.3,
               [[(c.S7B_MIX_LABEL, d.SANS, d.LABEL_SIZE, d.GRAY)]])
    d.add_text(s, 7.6, 2.8, 5.13, 1.3,
               [("mix", c.S7B_MIX_1, d.BODY_SIZE, d.INK),
                ("mix", c.S7B_MIX_2, d.BODY_SIZE, d.INK),
                ("mix", c.S7B_SPEARMAN, d.LABEL_SIZE, d.GRAY)],
               space_after=6)

    d.add_text(s, 0.6, 5.75, d.CONTENT_W, 0.45,
               [[(c.S7B_LINE, d.SANS_SB, 16, d.ORANGE)]])
    d.add_notes(s, c.S7B_NOTES)


# --- slide 8 --------------------------------------------------------------

def slide_8(prs):
    s = new_slide(prs, footer="8")
    d.add_title(s, c.S8_TITLE)

    d.stroke_box(s, 0.6, 2.6, 10.4, 1.55, fill=d.TINT5, stroke=None)
    d.add_text(s, 0.9, 2.8, 9.8, 1.15,
               [[(c.S8_REC_LABEL, d.SANS, d.LABEL_SIZE, d.GRAY)],
                [(c.S8_REC, d.SANS_SB, 18, d.INK)]],
               space_after=6)
    d.stroke_box(s, 0.6, 4.4, 10.4, 1.85, fill=d.TINT5, stroke=None)
    box = d.add_text(s, 0.9, 4.6, 9.8, 1.45,
                     [[(c.S8_PILOT_LABEL, d.SANS, d.LABEL_SIZE, d.GRAY)]],
                     space_after=6)
    p = box.text_frame.add_paragraph()
    p.line_spacing = 1.2
    r = p.add_run()
    r.text = c.S8_PILOT_LEAD
    d._style(r, d.SANS_SB, 15, d.INK)
    d.fill_runs(p, c.S8_PILOT_A, d.SANS, d.MONO, 15, d.INK)
    r = p.add_run()
    r.text = c.S8_PILOT_THRESH
    d._style(r, d.MONO, 15, d.ORANGE)
    d.fill_runs(p, c.S8_PILOT_B, d.SANS, d.MONO, 15, d.INK)
    d.add_notes(s, c.S8_NOTES)


# --- slide 9 --------------------------------------------------------------

def slide_9(prs):
    s = new_slide(prs, footer="9")
    d.add_title(s, c.S9_TITLE)

    y = 2.15
    for i, item in enumerate(c.S9_ITEMS, start=1):
        d.add_text(s, d.MARGIN, y, 0.5, 0.45,
                   [[(str(i), d.MONO, 20, d.GRAY)]],
                   wrap=False, space_after=0)
        box = d.add_text(s, 1.25, y + 0.02, 11.48, 0.9, [],
                         space_after=0, line_spacing=1.15)
        p = box.text_frame.paragraphs[0]
        _fill_limitation(p, item)
        y += 0.62 if len(item) < 60 else 0.95
    d.add_notes(s, c.S9_NOTES)


def _fill_limitation(p, text):
    """14pt ink; the WHERE clause in mono; '14 of 18' orange."""
    segments = [(text, d.INK)]
    if c.S9_ORANGE_SPAN in text:
        pre, post = text.split(c.S9_ORANGE_SPAN, 1)
        segments = [(pre, d.INK), (c.S9_ORANGE_SPAN, d.ORANGE),
                    (post, d.INK)]
    code = "WHERE new_product_id IS NOT NULL"
    for seg_text, color in segments:
        if code in seg_text:
            pre, post = seg_text.split(code, 1)
            parts = [(pre, d.SANS, color), (code, d.MONO, color),
                     (post, d.SANS, color)]
        else:
            parts = [(seg_text, d.SANS, color)]
        for t, font, col in parts:
            if not t:
                continue
            d.fill_runs(p, t, font, d.MONO if font == d.SANS else font,
                        d.BODY_SIZE, col)


# --- slide 10 -------------------------------------------------------------

def slide_10(prs):
    s = new_slide(prs, footer="10")
    d.add_title(s, c.S10_TITLE)
    d.add_text(s, d.MARGIN, 3.3, d.CONTENT_W, 1.1,
               [[(c.S10_AMOUNT, d.MONO, 54, d.ORANGE)]],
               wrap=False, space_after=0, line_spacing=1.0)
    d.add_text(s, d.MARGIN, 4.5, 10.5, 0.6,
               [("mix", c.S10_LINE, 16, d.INK)],
               space_after=0)
    d.add_notes(s, c.S10_NOTES)


# --- appendix -------------------------------------------------------------

def slide_a1(prs):
    s = new_slide(prs, footer="A1")
    d.add_title(s, c.A1_TITLE)
    cells = [c.A1_HEADER]
    for comp, val, wt, contrib in c.A1_ROWS:
        color = d.ORANGE if "0.668" in contrib else d.INK
        cells.append([comp,
                      [(val, d.MONO, d.BODY_SIZE, d.INK)],
                      [(wt, d.MONO, d.BODY_SIZE, d.INK)],
                      [(contrib, d.MONO, d.BODY_SIZE, color)]])
    d.tufte_table(s, 0.6, 2.3, [3.4, 1.5, 1.5, 2.4],
                  [0.38] + [0.52] * 4, cells,
                  align=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT,
                         PP_ALIGN.RIGHT],
                  last_row_totals=True)
    d.add_text(s, 0.6, 5.3, d.CONTENT_W, 0.6,
               [[(c.A1_FOOTNOTE, d.SANS, d.FOOT_SIZE, d.GRAY)]])
    d.add_notes(s, c.A1_NOTES)


def slide_a2(prs):
    s = new_slide(prs, footer="A2")
    d.add_title(s, c.A2_TITLE)
    _threshold_table(s, 0.6, 2.3, c.A2_HEADER, c.A2_ROWS, orange_row=1)
    d.add_notes(s, c.A2_NOTES)


def slide_a3(prs):
    """Executed poison-fixture capture (item 24). Requires the artifacts
    from `python deck/poison_fixture.py`; per the hard gate, A3 is dropped
    (not shipped as a placeholder) when they are missing."""
    meta_path = HERE / "out" / "fixture" / "poison_fixture_meta.json"
    png_path = HERE / "out" / "fixture" / "poison_fixture_log.png"
    if not (meta_path.exists() and png_path.exists()):
        print("WARNING: poison-fixture artifacts missing; A3 dropped")
        return
    meta = json.loads(meta_path.read_text(encoding="utf-8"))
    if not meta.get("gate_fired"):
        raise SystemExit("poison-fixture gate did NOT fire; stop the build")

    s = new_slide(prs, footer="A3")
    d.add_title(s, c.A3_TITLE)
    d.add_text(s, d.MARGIN, 1.95, d.CONTENT_W, 0.35,
               [[(c.A3_GATE_LINE, d.SANS_SB, 14, d.ORANGE)]])
    from PIL import Image
    with Image.open(png_path) as img:
        aspect = img.height / img.width
    pic_w = min(11.2, 4.1 / aspect)
    s.shapes.add_picture(str(png_path), d.Inches(d.MARGIN), d.Inches(2.45),
                         d.Inches(pic_w), d.Inches(pic_w * aspect))
    d.add_text(s, d.MARGIN, 2.65 + pic_w * aspect, d.CONTENT_W, 0.55,
               [[(c.A3_CAPTION.format(**meta), d.SANS, d.FOOT_SIZE, d.GRAY)]])
    d.add_notes(s, c.A3_NOTES)


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = d.SLIDE_W_EMU
    prs.slide_height = d.SLIDE_H_EMU
    cover(prs)
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7a(prs)
    slide_7b(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    slide_a1(prs)
    slide_a2(prs)
    slide_a3(prs)
    path = OUT / "BZAN545_Final_Deck.pptx"
    prs.save(path)
    print(f"wrote {path} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
