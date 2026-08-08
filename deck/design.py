"""
Design system for the BZAN 545 final deck. Binding spec lives in the build
prompt; plan.md at the worktree root is the sole source of slide content.

Canvas: 16:9 (12192000 x 6858000 EMU), off-white F7F6F3, ink 1A1A1A,
margins 0.6in. Titles IBM Plex Sans SemiBold 24pt top-left, identical x/y on
every content slide. Numerals in tables/callouts: IBM Plex Mono. UT orange
FF8200 marks exactly one element per slide.

Font names must be the GDI family names from the installed TTFs' name-table
ID 1 ("IBM Plex Sans SmBld", not "IBM Plex Sans SemiBold") or PowerPoint
silently substitutes.
"""

import re

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# --- canvas ---------------------------------------------------------------
SLIDE_W_EMU = 12192000          # 13.333 in
SLIDE_H_EMU = 6858000           # 7.5 in
SLIDE_W_IN = 13.3333
SLIDE_H_IN = 7.5
MARGIN = 0.6
CONTENT_W = SLIDE_W_IN - 2 * MARGIN     # 12.133 in

TITLE_X = MARGIN
TITLE_Y = 0.55
TITLE_W = CONTENT_W
TITLE_H = 1.05                  # two lines at 24pt

# --- color ----------------------------------------------------------------
BG = RGBColor.from_string("F7F6F3")
INK = RGBColor.from_string("1A1A1A")
GRAY = RGBColor.from_string("666666")
HAIRLINE = RGBColor.from_string("D9D6D0")
ORANGE = RGBColor.from_string("FF8200")
TINT5 = RGBColor.from_string("ECEBE8")   # 5% ink over bg
TINT8 = RGBColor.from_string("E5E4E2")   # 8% ink over bg

# --- type (GDI family names; see module docstring) ------------------------
SANS = "IBM Plex Sans"
SANS_SB = "IBM Plex Sans SmBld"
SANS_MD = "IBM Plex Sans Medm"
MONO = "IBM Plex Mono"
MONO_SB = "IBM Plex Mono SmBld"

TITLE_SIZE = 24
BODY_SIZE = 14
LABEL_SIZE = 12
FOOT_SIZE = 10

# numeric tokens get IBM Plex Mono inside mixed text
NUM_RE = re.compile(r"\$?\d[\d,.\-]*\d%?|\$?\d%?")


# --- slide chrome ---------------------------------------------------------

def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_footer(slide, label):
    """Slide number only, 10pt gray, bottom-right."""
    box = slide.shapes.add_textbox(
        Inches(SLIDE_W_IN - MARGIN - 0.8), Inches(SLIDE_H_IN - 0.42),
        Inches(0.8), Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(label)
    _style(r, MONO, FOOT_SIZE, GRAY)
    return box


def add_title(slide, text):
    """Assertion title, verbatim, sentence case, top-left, never centered."""
    box = slide.shapes.add_textbox(
        Inches(TITLE_X), Inches(TITLE_Y), Inches(TITLE_W), Inches(TITLE_H))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.05
    r = p.add_run()
    r.text = text
    _style(r, SANS_SB, TITLE_SIZE, INK)
    return box


def add_notes(slide, lines):
    """Speaker notes, one paragraph per entry (no bullet characters)."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line


# --- text -----------------------------------------------------------------

def _style(run, font, size, color, italic=False):
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.color.rgb = color
    f.italic = italic
    # keep east-asian/cs in sync so nothing falls back to a theme font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)


def fill_runs(p, text, sans, mono, size, color):
    """Write text into paragraph p, numeric tokens in mono, rest in sans."""
    pos = 0
    for m in NUM_RE.finditer(text):
        if m.start() > pos:
            r = p.add_run()
            r.text = text[pos:m.start()]
            _style(r, sans, size, color)
        r = p.add_run()
        r.text = m.group()
        _style(r, mono, size, color)
        pos = m.end()
    if pos < len(text):
        r = p.add_run()
        r.text = text[pos:]
        _style(r, sans, size, color)


def add_text(slide, x, y, w, h, entries, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True, space_after=6,
             line_spacing=1.15):
    """entries: list of paragraphs; each is a list of
    (text, font, size, color) run tuples, or ("mix", text, size, color)
    to auto-set numerals in mono."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, runs in enumerate(entries):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if runs and runs[0] == "mix":
            _, text, size, color = runs
            fill_runs(p, text, SANS, MONO, size, color)
        else:
            for text, font, size, color in runs:
                r = p.add_run()
                r.text = text
                _style(r, font, size, color)
    return box


# --- lines and shapes -----------------------------------------------------

def hairline(slide, x, y, w, weight=0.75, color=HAIRLINE):
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def stroke_box(slide, x, y, w, h, fill=None, stroke=INK, weight=1.0):
    """Square-corner rectangle; fill None = transparent, or a tint."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        shp.line.width = Pt(weight)
    return shp


def connector(slide, x1, y1, x2, y2, elbow=False, arrow=True, weight=1.0,
              color=INK):
    kind = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    ln = slide.shapes.add_connector(
        kind, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    if arrow:
        lnEl = ln.line._get_or_add_ln()
        tail = lnEl.makeelement(qn("a:tailEnd"),
                                {"type": "triangle", "w": "sm", "len": "sm"})
        lnEl.append(tail)
    return ln


def incident_dot(slide, cx, cy, d, missed=False):
    """Incident state glyph, identical geometry at every scale.
    caught = filled ink dot; missed = open dot stroked orange (bg fill)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
        Inches(d), Inches(d))
    shp.shadow.inherit = False
    if missed:
        shp.fill.solid()
        shp.fill.fore_color.rgb = BG
        shp.line.color.rgb = ORANGE
        shp.line.width = Pt(1.0)
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = INK
        shp.line.fill.background()
    return shp


def incident_strip(slide, x, y, w, dot_d, missed_index, dates=None,
                   label_size=LABEL_SIZE):
    """Six-dot incident strip (five filled, one open-orange), evenly spaced
    across w with a hairline behind. Optional date labels beneath each dot.
    Same routine draws the full-width slide-1 strip and the slide-4 corner
    recurrence, so geometry stays identical."""
    n = 6
    hairline(slide, x, y, w)
    step = w / (n - 1)
    for i in range(n):
        cx = x + i * step
        incident_dot(slide, cx, y, dot_d, missed=(i == missed_index))
        if dates:
            add_text(slide, cx - step / 2, y + dot_d / 2 + 0.08, step, 0.3,
                     [[(dates[i], MONO, label_size, GRAY)]],
                     align=PP_ALIGN.CENTER, wrap=False, space_after=0)


# --- Tufte table ----------------------------------------------------------

def _cell_border(tc, edge, color=None, weight=None):
    """edge in lnL/lnR/lnT/lnB; color None -> explicit no line."""
    tcPr = tc.get_or_add_tcPr()
    tag = qn("a:" + edge)
    for e in tcPr.findall(tag):
        tcPr.remove(e)
    ln = tcPr.makeelement(tag, {})
    if color is None:
        ln.append(ln.makeelement(qn("a:noFill"), {}))
    else:
        ln.set("w", str(int(weight * 12700)))
        ln.set("cap", "flat")
        fillEl = ln.makeelement(qn("a:solidFill"), {})
        clr = fillEl.makeelement(qn("a:srgbClr"), {"val": str(color)})
        fillEl.append(clr)
        ln.append(fillEl)
    # schema order: lnL, lnR, lnT, lnB first inside tcPr
    order = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
    idx = order.index("a:" + edge)
    before = [qn(t) for t in order[idx + 1:]]
    inserted = False
    for child in tcPr:
        if child.tag in before or child.tag not in [qn(t) for t in order]:
            child.addprevious(ln)
            inserted = True
            break
    if not inserted:
        tcPr.append(ln)


def tufte_table(slide, x, y, col_widths, row_heights, cells,
                header_size=LABEL_SIZE, body_size=BODY_SIZE,
                align=None, last_row_totals=False):
    """No vertical rules, no fills, no banding. 0.75pt hairline under header
    (and above totals row if last_row_totals). Header 12pt semibold.
    cells[r][c] is a string ("mix" numeral treatment) or list of run tuples.
    align: optional list of PP_ALIGN per column."""
    n_rows, n_cols = len(cells), len(col_widths)
    gf = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y),
        Inches(sum(col_widths)), Inches(sum(row_heights)))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    # drop the theme table style entirely
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        for style_id in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(style_id)
    for c, wid in enumerate(col_widths):
        tbl.columns[c].width = Inches(wid)
    for r, hgt in enumerate(row_heights):
        tbl.rows[r].height = Inches(hgt)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            cell.fill.background()
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0 if c == 0 else 0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            tc = cell._tc
            for edge in ("lnL", "lnR", "lnT", "lnB"):
                _cell_border(tc, edge)
            if r == 0:
                _cell_border(tc, "lnB", "D9D6D0", 0.75)
            # PowerPoint resolves inner horizontal borders from the upper
            # cell's lnB, so the totals rule goes on the row above
            if last_row_totals and r == n_rows - 2:
                _cell_border(tc, "lnB", "D9D6D0", 0.75)
            content = cells[r][c]
            p = cell.text_frame.paragraphs[0]
            p.alignment = (align[c] if align else PP_ALIGN.LEFT)
            if r == 0:
                runs = [(content, SANS_SB, header_size, INK)] \
                    if isinstance(content, str) else content
            elif isinstance(content, str):
                fill_runs(p, content, SANS, MONO, body_size, INK)
                continue
            else:
                runs = content
            for text, font, size, color in runs:
                run = p.add_run()
                run.text = text
                _style(run, font, size, color)
    return gf


# --- display figures ------------------------------------------------------

def stat(slide, x, y, w, value, label, size=44, color=INK,
         label_color=GRAY, align=PP_ALIGN.LEFT, mono=True):
    """Key-figure display treatment: 40-60pt mono value, 12pt label beneath."""
    add_text(slide, x, y, w, size / 72 + 0.15,
             [[(value, MONO if mono else SANS_SB, size, color)]],
             align=align, wrap=False, space_after=0, line_spacing=1.0)
    add_text(slide, x, y + size / 72 + 0.12, w, 0.55,
             [[(label, SANS, LABEL_SIZE, label_color)]],
             align=align, space_after=0)
